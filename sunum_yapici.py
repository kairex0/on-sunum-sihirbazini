import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from fpdf import FPDF
import matplotlib.pyplot as plt
import io

# -------------------------------------------------------------
# TÜRKÇE KARAKTER TEMİZLEME FONKSİYONU
# -------------------------------------------------------------
def turkce_temizle(metin):
    kaynak = "şŞıİğĞüÜöÖçÇ "
    hedef  = "sSiiGGuuooCC_" # Boşlukları da alt tire yapıyoruz ki dosya adında şık dursun
    tablo = str.maketrans(kaynak, hedef)
    return metin.translate(tablo)

# -------------------------------------------------------------
# 1. POWERPOINT MOTORU
# -------------------------------------------------------------
def sablonlu_sunum_uret(ilce, mahalle, m2, imar, arsa_fiyati, danisman, konsept, sablon_turu, yuklenen_resim):
    prs = Presentation()
    
    if sablon_turu == "ON Premium (Gold & Lacivert)":
        ANA_RENK = RGBColor(11, 29, 58)      
        VURGU_RENK = RGBColor(212, 175, 55)   
        METIN_KOYU = RGBColor(40, 45, 55)
        GRAFIK_RENKLER = ['#0B1D3A', '#D4AF37']
    elif sablon_turu == "ON Nature (Doğa & Toprak)":
        ANA_RENK = RGBColor(34, 76, 56)       
        VURGU_RENK = RGBColor(194, 143, 83)   
        METIN_KOYU = RGBColor(50, 60, 50)
        GRAFIK_RENKLER = ['#224C38', '#C28F53']
    else: 
        ANA_RENK = RGBColor(43, 43, 43)       
        VURGU_RENK = RGBColor(241, 90, 36)    
        METIN_KOYU = RGBColor(30, 30, 30)
        GRAFIK_RENKLER = ['#2B2B2B', '#F15A24']
        
    BEYAZ = RGBColor(255, 255, 255)
    
    # SLAYT 1: KAPAK
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = ANA_RENK; bg.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"İZMİR {ilce.upper()} - {mahalle.upper()}"
    p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = VURGU_RENK
    
    p2 = tf.add_paragraph()
    p2.text = f"{konsept.upper()} PROJE GELİŞTİRME & YATIRIM ANALİZİ"
    p2.font.size = Pt(18); p2.font.color.rgb = BEYAZ; p2.space_before = Pt(10)
    
    p3 = tf.add_paragraph()
    p3.text = f"\nHazırlayan: {danisman} | ON Türkiye"
    p3.font.size = Pt(14); p3.font.color.rgb = VURGU_RENK

    # SLAYT 2: PORTFÖY ÖZELLİKLERİ
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    t_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    t_box.text_frame.paragraphs[0].text = "PORTFÖY ÖZELLİKLERİ VE DETAYLAR"
    t_box.text_frame.paragraphs[0].font.size = Pt(26); t_box.text_frame.paragraphs[0].font.bold = True; t_box.text_frame.paragraphs[0].font.color.rgb = ANA_RENK
    
    tf2 = slide2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5), Inches(5)).text_frame
    maddeler = [
        f"📍 Konum: İzmir / {ilce} / {mahalle}",
        f"📐 Alan: {m2} m²",
        f"📜 İmar: {imar}",
        "⚡ Altyapı: Elektrik & Su hatları hazır.",
        f"💰 Alım Bedeli: {arsa_fiyati:,.0f} TL".replace(",", ".")
    ]
    for m in maddeler:
        p_m = tf2.add_paragraph(); p_m.text = m; p_m.font.size = Pt(16); p_m.font.color.rgb = METIN_KOYU; p_m.space_after = Pt(12)
        
    if yuklenen_resim is not None:
        image_stream = io.BytesIO(yuklenen_resim.getvalue())
        slide2.shapes.add_picture(image_stream, Inches(5.8), Inches(1.5), width=Inches(3.8))

    # SLAYT 3: MALİYET VE GRAFİK SLAYTI
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    t_box3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    t_box3.text_frame.paragraphs[0].text = "YATIRIM BÜTÇESİ VE DAĞILIMI"
    t_box3.text_frame.paragraphs[0].font.size = Pt(26); t_box3.text_frame.paragraphs[0].font.bold = True; t_box3.text_frame.paragraphs[0].font.color.rgb = ANA_RENK
    
    tf3 = slide3.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5), Inches(5)).text_frame
    
    taban_alan = int(int(m2) * 0.15)
    toplam_insaat = taban_alan * 2
    insaat_maliyeti = toplam_insaat * 30000
    toplam_maliyet = arsa_fiyati + insaat_maliyeti
    tahmini_bitis_degeri = toplam_maliyet * 1.6
    
    finansallar = [
        f"🏗️ İnşaat Alanı: {toplam_insaat} m²",
        f"💵 Arsa Bedeli: {arsa_fiyati:,.0f} TL".replace(",", "."),
        f"🔨 Yapım Maliyeti: {insaat_maliyeti:,.0f} TL".replace(",", "."),
        f"📊 Toplam Bütçe: {toplam_maliyet:,.0f} TL".replace(",", "."),
        f"📈 Bitiş Değeri: {tahmini_bitis_degeri:,.0f} TL".replace(",", ".")
    ]
    for f in finansallar:
        p_f = tf3.add_paragraph(); p_f.text = f; p_f.font.size = Pt(15); p_f.font.color.rgb = METIN_KOYU; p_f.space_after = Pt(12)

    labels = ['Arsa Bedeli', 'Insaat Maliyeti']
    sizes = [arsa_fiyati, insaat_maliyeti]
    
    fig, ax = plt.subplots(figsize=(4, 4))
    ax.pie(sizes, labels=labels, autopct='%1.1f%%', startangle=90, colors=GRAFIK_RENKLER, textprops={'fontsize': 12, 'weight': 'bold'})
    ax.axis('equal')  
    
    chart_stream = io.BytesIO()
    plt.savefig(chart_stream, format='png', bbox_inches='tight', transparent=True)
    chart_stream.seek(0)
    plt.close()
    
    slide3.shapes.add_picture(chart_stream, Inches(5.5), Inches(1.5), width=Inches(4.0))

    binary_output = io.BytesIO()
    prs.save(binary_output)
    binary_output.seek(0)
    return binary_output

# -------------------------------------------------------------
# 2. GÜVENLİ PDF ÜRETME MOTORU
# -------------------------------------------------------------
def pdf_rapor_uret(ilce, mahalle, m2, imar, arsa_fiyati, danisman, konsept):
    pdf = FPDF()
    pdf.add_page()
    
    pdf.set_font("Courier", "B", 20) 
    pdf.set_text_color(11, 29, 58) 
    pdf.cell(0, 15, turkce_temizle("ON TURKIYE GAYRIMENKUL"), ln=True, align="C")
    
    pdf.set_font("Courier", "", 12)
    pdf.cell(0, 10, turkce_temizle("YATIRIM VE PROJE ANALIZ RAPORU"), ln=True, align="C")
    pdf.line(10, 40, 200, 40)
    pdf.ln(10)
    
    pdf.set_font("Courier", "B", 14)
    pdf.set_text_color(241, 90, 36) 
    portfoy_baslik = f"PORTFOY: IZMIR {ilce.upper()} - {mahalle.upper()}"
    pdf.cell(0, 10, turkce_temizle(portfoy_baslik), ln=True)
    pdf.ln(5)
    
    pdf.set_font("Courier", "", 11)
    pdf.set_text_color(50, 50, 50)
    
    fiyat_str = f"{arsa_fiyati:,.0f} TL".replace(",", ".")
    taban_alan = int(int(m2) * 0.15)
    toplam_insaat = taban_alan * 2
    insaat_maliyeti = toplam_insaat * 30000
    toplam_maliyet = arsa_fiyati + insaat_maliyeti
    tahmini_bitis_degeri = toplam_maliyet * 1.6
    
    pdf.cell(0, 8, turkce_temizle(f"- Toplam Arazi Alani: {m2} m2"), ln=True)
    pdf.cell(0, 8, turkce_temizle(f"- Imar Durumu: {imar}"), ln=True)
    pdf.cell(0, 8, turkce_temizle(f"- Onerilen Proje Konsepti: {konsept}"), ln=True)
    pdf.cell(0, 8, turkce_temizle(f"- Arsa Alim Bedeli: {fiyat_str}"), ln=True)
    pdf.ln(5)
    
    pdf.set_font("Courier", "B", 13)
    pdf.set_text_color(11, 29, 58)
    pdf.cell(0, 10, turkce_temizle("PROJEKSIYON VE MALIYET TABLOSU"), ln=True)
    pdf.set_font("Courier", "", 11)
    
    pdf.cell(0, 8, turkce_temizle(f"- Planlanan Insaat Alani: {toplam_insaat} m2"), ln=True)
    pdf.cell(0, 8, turkce_temizle(f"- Tahmini Yapim Maliyeti: {insaat_maliyeti:,.0f} TL".replace(",", ".")), ln=True)
    pdf.cell(0, 8, turkce_temizle(f"- Toplam Yatirim Butcesi (Arsa+Insaat): {toplam_maliyet:,.0f} TL".replace(",", ".")), ln=True)
    pdf.cell(0, 8, turkce_temizle(f"- Proje Bitisindeki Tahmini Deger: {tahmini_bitis_degeri:,.0f} TL".replace(",", ".")), ln=True)
    
    pdf.ln(15)
    pdf.line(10, pdf.get_y(), 200, pdf.get_y())
    pdf.ln(5)
    pdf.set_font("Courier", "I", 10)
    rapor_hazirlayan = f"Raporu Hazirlayan: {danisman} - ON Turkiye"
    pdf.cell(0, 10, turkce_temizle(rapor_hazirlayan), align="R")
    
    return bytes(pdf.output())

# -------------------------------------------------------------
# 3. DİNAMİK İLAN METNİ MOTORU
# -------------------------------------------------------------
def ilan_metni_uret(ilce, mahalle, m2, imar, arsa_fiyati, konsept, danisman):
    fiyat_str = f"{arsa_fiyati:,.0f} TL".replace(",", ".")
    metin = f"""🚀 ON TÜRKİYE'DEN İZMİR {ilce.upper()} {mahalle.upper()}'DE KAÇIRILMAYACAK YATIRIM FIRSATI!\n\n✨ Portföy Özellikleri:\n• Konum: İzmir / {ilce} / {mahalle}\n• Genişlik: {m2} m²\n• İmar Durumu: {imar}\n• Altyapı: Elektrik ve su hazır!\n\n🎯 Önerilen Proje: {konsept}\n💵 Fırsat Alım Bedeli: {fiyat_str}\n\nYatırım Danışmanı: {danisman}\n🏢 ON TÜRKİYE GAYRİMENKUL"""
    return metin

# -------------------------------------------------------------
# STREAMLIT ARAYÜZÜ
# -------------------------------------------------------------
st.set_page_config(page_title="ON Türkiye Sunum Sihirbazı v6", page_icon="🏢", layout="centered")

st.title("🏢 ON Türkiye Sunum Sihirbazı v6")
st.write("Bilgileri girin; dinamik dosya ismiyle raporlarınızı indirin!")
st.divider()

col1, col2 = st.columns(2)
with col1:
    ilce = st.text_input("İlçe Adı", value="Urla")
    m2 = st.number_input("Metrekare (m²)", min_value=1, value=500)
    danisman = st.text_input("Danışman Adı Soyadı", value="Fatih Türkdönmez")
with col2:
    mahalle = st.text_input("Mahalle / Köy", value="Kuşçular")
    imar = st.text_input("İmar Durumu", value="%15/30 Konut İmarlı")
    arsafiyati = st.number_input("Arsa Fiyatı (TL)", min_value=0, value=6500000, step=50000)

st.divider()
sablon_turu = st.selectbox("Sunum Şablonu Tasarımı (Renk Modu)", ["ON Premium (Gold & Lacivert)", "ON Nature (Doğa & Toprak)", "ON Commercial (Modern & Antrasit)"])
konsept = st.selectbox("Önerilecek Proje Konsepti", ["Premium Taş Ev", "Eko-Tiny House Yaşam Alanı", "Yüksek Getirili Villa"])
yuklenen_resim = st.file_uploader("Arsa / Drone Fotoğrafı Yükleyin (Opsiyonel)", type=["jpg", "jpeg", "png"])
st.divider()

# 🏷️ DİNAMİK DOSYA İSMİ GENERATORU
temiz_ilce = turkce_temizle(ilce.lower())
temiz_mahalle = turkce_temizle(mahalle.lower())
temiz_konsept = turkce_temizle(konsept.lower())

pptx_dosya_adi = f"ON_{temiz_ilce}_{temiz_mahalle}_{temiz_konsept}_sunumu.pptx"
pdf_dosya_adi = f"ON_{temiz_ilce}_{temiz_mahalle}_{temiz_konsept}_analiz_raporu.pdf"

btn_col1, btn_col2, btn_col3 = st.columns(3)

with btn_col1:
    if st.button("🚀 PPTX Sunumu Üret", use_container_width=True):
        sunum_dosyasi = sablonlu_sunum_uret(ilce, mahalle, str(m2), imar, arsafiyati, danisman, konsept, sablon_turu, yuklenen_resim)
        st.download_button(label="📥 PowerPoint İndir", data=sunum_dosyasi, file_name=pptx_dosya_adi, mime="application/vnd.openxmlformats-officedocument.presentationml.presentation", use_container_width=True)

with btn_col2:
    if st.button("📄 Yatırım PDF Raporu Bas", use_container_width=True):
        pdf_dosyasi = pdf_rapor_uret(ilce, mahalle, str(m2), imar, arsafiyati, danisman, konsept)
        st.download_button(label="📥 PDF Raporunu İndir", data=pdf_dosyasi, file_name=pdf_dosya_adi, mime="application/pdf", use_container_width=True)

with btn_col3:
    if st.button("✍️ İlan Metni Üret", use_container_width=True):
        st.session_state["ilan_metni"] = ilan_metni_uret(ilce, mahalle, str(m2), imar, arsafiyati, konsept, danisman)

if "ilan_metni" in st.session_state:
    st.text_area(label="Hazır İlan Açıklaması:", value=st.session_state["ilan_metni"], height=200)